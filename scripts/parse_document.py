#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""Parse user-uploaded documents into page/block JSON for the PPT agent."""

from __future__ import annotations

import json
import os
import re
import sys
from html import unescape
from pathlib import Path
from typing import Any, Dict, List


MAX_BLOCK_CHARS = 520
MAX_TOTAL_BLOCKS = 120


def clean_text(value: Any) -> str:
    text = "" if value is None else str(value)
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def short(value: str, limit: int = MAX_BLOCK_CHARS) -> str:
    value = clean_text(value)
    if len(value) <= limit:
        return value
    return value[: limit - 1].rstrip() + "…"


def block_type(text: str, fallback: str = "text") -> str:
    stripped = clean_text(text)
    if not stripped:
        return fallback
    if len(stripped) <= 32 and re.search(r"(第[一二三四五六七八九十0-9]+章|摘要|目录|背景|目标|方案|结论|建议|Overview|Summary)", stripped, re.I):
        return "heading"
    if re.match(r"^([•\-*]|[0-9一二三四五六七八九十]+[.、])\s+", stripped):
        return "list"
    if stripped.count("|") >= 2 or "\t" in stripped:
        return "table"
    if len(stripped) <= 26:
        return "heading"
    return fallback


def add_page(pages: List[Dict[str, Any]], blocks: List[Dict[str, Any]], page_no: int) -> None:
    page_blocks = [item for item in blocks if item["page"] == page_no]
    if not page_blocks:
        return
    title = next((item["text"] for item in page_blocks if item["type"] in {"title", "heading"}), page_blocks[0]["text"])
    summary = "；".join(item["text"] for item in page_blocks[:3])
    pages.append(
        {
            "page": page_no,
            "title": short(title, 42),
            "summary": short(summary, 180),
            "blockCount": len(page_blocks),
            "imageCount": sum(1 for item in page_blocks if item["type"] == "image"),
            "tableCount": sum(1 for item in page_blocks if item["type"] == "table"),
            "blocks": page_blocks[:12],
        }
    )


def make_block(file_stem: str, page: int, idx: int, text: str, kind: str = "text", confidence: int = 88) -> Dict[str, Any]:
    return {
        "id": f"{file_stem}-p{page}-b{idx}",
        "page": page,
        "type": kind,
        "text": short(text),
        "confidence": confidence,
        "sourceRef": f"p{page}/b{idx}",
    }


def split_paragraphs(text: str) -> List[str]:
    text = clean_text(text)
    if not text:
        return []
    parts = re.split(r"\n\s*\n|(?<=[。！？!?])\s+(?=[一-龥A-Za-z0-9])", text)
    cleaned = [short(part, MAX_BLOCK_CHARS) for part in parts if len(clean_text(part)) >= 8]
    if len(cleaned) <= 1 and len(text) > MAX_BLOCK_CHARS:
        cleaned = [short(text[i : i + MAX_BLOCK_CHARS], MAX_BLOCK_CHARS) for i in range(0, len(text), MAX_BLOCK_CHARS)]
    return cleaned[:18]


def parse_pdf(path: Path, stem: str) -> Dict[str, Any]:
    import pdfplumber

    blocks: List[Dict[str, Any]] = []
    pages: List[Dict[str, Any]] = []
    with pdfplumber.open(str(path)) as pdf:
        for page_index, page in enumerate(pdf.pages, 1):
            page_block_idx = 1
            text = clean_text(page.extract_text(x_tolerance=1, y_tolerance=3) or "")
            for paragraph in split_paragraphs(text):
                blocks.append(make_block(stem, page_index, page_block_idx, paragraph, block_type(paragraph), 88))
                page_block_idx += 1
            try:
                tables = page.extract_tables() or []
            except Exception:
                tables = []
            for table in tables[:3]:
                rows = []
                for row in table[:8]:
                    cells = [clean_text(cell) for cell in row if clean_text(cell)]
                    if cells:
                        rows.append(" | ".join(cells))
                if rows:
                    blocks.append(make_block(stem, page_index, page_block_idx, "\n".join(rows), "table", 82))
                    page_block_idx += 1
            try:
                images = getattr(page, "images", []) or []
            except Exception:
                images = []
            for image_index, _ in enumerate(images[:4], 1):
                blocks.append(make_block(stem, page_index, page_block_idx, f"图片素材 {image_index}：来自 PDF 第 {page_index} 页，需在后续视觉设计中引用或重绘。", "image", 72))
                page_block_idx += 1
            add_page(pages, blocks, page_index)
    return {"pages": pages, "blocks": blocks, "sourceKind": "pdf"}


def parse_docx(path: Path, stem: str) -> Dict[str, Any]:
    from docx import Document

    document = Document(str(path))
    blocks: List[Dict[str, Any]] = []
    idx = 1
    for paragraph in document.paragraphs:
        text = clean_text(paragraph.text)
        if len(text) < 2:
            continue
        style_name = getattr(paragraph.style, "name", "") or ""
        kind = "heading" if "Heading" in style_name or "标题" in style_name or len(text) <= 28 else block_type(text)
        blocks.append(make_block(stem, 1, idx, text, kind, 90))
        idx += 1
    for table_index, table in enumerate(document.tables, 1):
        rows = []
        for row in table.rows[:12]:
            cells = [clean_text(cell.text) for cell in row.cells if clean_text(cell.text)]
            if cells:
                rows.append(" | ".join(cells))
        if rows:
            blocks.append(make_block(stem, 1, idx, f"表格 {table_index}\n" + "\n".join(rows), "table", 84))
            idx += 1
    pages: List[Dict[str, Any]] = []
    add_page(pages, blocks, 1)
    return {"pages": pages, "blocks": blocks, "sourceKind": "docx"}


def parse_pptx(path: Path, stem: str) -> Dict[str, Any]:
    from pptx import Presentation

    presentation = Presentation(str(path))
    blocks: List[Dict[str, Any]] = []
    pages: List[Dict[str, Any]] = []
    for page_no, slide in enumerate(presentation.slides, 1):
        idx = 1
        for shape in slide.shapes:
            text = clean_text(getattr(shape, "text", ""))
            if not text:
                continue
            kind = "title" if idx == 1 and len(text) <= 48 else block_type(text)
            blocks.append(make_block(stem, page_no, idx, text, kind, 88))
            idx += 1
        add_page(pages, blocks, page_no)
    return {"pages": pages, "blocks": blocks, "sourceKind": "pptx"}


def decode_text_bytes(data: bytes) -> str:
    try:
        utf8_text = data.decode("utf-8-sig")
        if "\ufffd" not in utf8_text:
            return utf8_text.strip()
    except UnicodeDecodeError:
        pass

    def score(value: str) -> float:
        chinese = len(re.findall(r"[\u3400-\u9fff]", value))
        ascii_chars = len(re.findall(r"[A-Za-z0-9]", value))
        replacements = value.count("\ufffd")
        mojibake = len(re.findall(r"[锟�鏉窞佃涓]", value))
        punctuation = len(re.findall(r"[，。；：、！？]", value))
        return chinese * 4 + ascii_chars * 0.15 + punctuation * 0.8 - replacements * 12 - mojibake * 1.5

    candidates = []
    for encoding in ("utf-8-sig", "utf-8", "gb18030", "cp936"):
        try:
            candidates.append(data.decode(encoding))
        except UnicodeDecodeError:
            candidates.append(data.decode(encoding, errors="replace"))
    return max(candidates, key=score).strip()


def parse_text(path: Path, stem: str) -> Dict[str, Any]:
    content = decode_text_bytes(path.read_bytes())
    paragraphs = split_paragraphs(content)
    blocks = [make_block(stem, 1, index + 1, paragraph, block_type(paragraph), 88) for index, paragraph in enumerate(paragraphs)]
    pages: List[Dict[str, Any]] = []
    add_page(pages, blocks, 1)
    return {"pages": pages, "blocks": blocks, "sourceKind": "text"}


def parse_html(path: Path, stem: str) -> Dict[str, Any]:
    content = decode_text_bytes(path.read_bytes())
    content = re.sub(r"(?is)<script[^>]*>.*?</script>", " ", content)
    content = re.sub(r"(?is)<style[^>]*>.*?</style>", " ", content)
    content = re.sub(r"(?is)<nav[^>]*>.*?</nav>", " ", content)
    content = re.sub(r"(?is)<header[^>]*>.*?</header>", " ", content)
    content = re.sub(r"(?is)<footer[^>]*>.*?</footer>", " ", content)
    content = re.sub(r"(?i)<br\s*/?>", "\n", content)
    content = re.sub(r"(?i)</(p|div|section|article|li|h[1-6]|tr)>", "\n", content)
    text = clean_text(unescape(re.sub(r"<[^>]+>", " ", content)))
    paragraphs = split_paragraphs(text)
    blocks = [make_block(stem, 1, index + 1, paragraph, block_type(paragraph), 86) for index, paragraph in enumerate(paragraphs)]
    pages: List[Dict[str, Any]] = []
    add_page(pages, blocks, 1)
    return {"pages": pages, "blocks": blocks, "sourceKind": "text"}


def parse_image(path: Path, stem: str) -> Dict[str, Any]:
    blocks = [
        make_block(
            stem,
            1,
            1,
            f"图片资料：{path.name}。当前本地未安装 OCR，已作为视觉参考素材保留；后续可接入 OCR/image2 模型提取图中文字和图表含义。",
            "image",
            58,
        )
    ]
    pages: List[Dict[str, Any]] = []
    add_page(pages, blocks, 1)
    return {"pages": pages, "blocks": blocks, "sourceKind": "image"}


def build_summary(blocks: List[Dict[str, Any]]) -> str:
    texts = [item["text"] for item in blocks if item["type"] != "image"]
    if not texts:
        texts = [item["text"] for item in blocks]
    return short(" ".join(texts[:8]), 420)


def outline_suggestions(blocks: List[Dict[str, Any]]) -> List[str]:
    headings = [item["text"] for item in blocks if item["type"] in {"title", "heading"} and len(item["text"]) >= 3]
    if len(headings) < 5:
        for item in blocks:
            text = item["text"]
            if 8 <= len(text) <= 48 and text not in headings:
                headings.append(text)
            if len(headings) >= 8:
                break
    return [short(item, 36) for item in headings[:10]]


def parse(path: Path, display_name: str | None = None) -> Dict[str, Any]:
    file_name = display_name or path.name
    suffix = (Path(file_name).suffix or path.suffix).lower()
    stem_source = Path(file_name).stem
    stem = re.sub(r"[^a-zA-Z0-9_-]+", "-", stem_source).strip("-") or "upload"
    if suffix == ".pdf":
        parsed = parse_pdf(path, stem)
    elif suffix == ".docx":
        parsed = parse_docx(path, stem)
    elif suffix == ".pptx":
        parsed = parse_pptx(path, stem)
    elif suffix in {".txt", ".md"}:
        parsed = parse_text(path, stem)
    elif suffix in {".html", ".htm"}:
        parsed = parse_html(path, stem)
    elif suffix in {".png", ".jpg", ".jpeg", ".webp"}:
        parsed = parse_image(path, stem)
    else:
        parsed = {"pages": [], "blocks": [], "sourceKind": "unknown"}

    blocks = parsed["blocks"][:MAX_TOTAL_BLOCKS]
    pages = parsed["pages"]
    return {
        "fileName": file_name,
        "fileType": suffix.lstrip(".") or "unknown",
        "pageCount": len(pages),
        "blockCount": len(blocks),
        "summary": build_summary(blocks),
        "outlineSuggestions": outline_suggestions(blocks),
        "pages": pages,
        "blocks": blocks,
        "sourceKind": parsed["sourceKind"],
    }


def main() -> int:
    if len(sys.argv) < 2:
        print(json.dumps({"error": "missing path"}, ensure_ascii=True))
        return 2
    path = Path(sys.argv[1])
    display_name = sys.argv[2] if len(sys.argv) >= 3 else None
    if not path.exists():
        print(json.dumps({"error": "file not found"}, ensure_ascii=True))
        return 2
    try:
        result = parse(path, display_name)
    except Exception as exc:  # noqa: BLE001 - command-line helper returns machine-readable errors.
        result = {
            "fileName": display_name or path.name,
            "fileType": path.suffix.lower().lstrip(".") or "unknown",
            "pageCount": 0,
            "blockCount": 0,
            "summary": f"解析失败：{exc}",
            "outlineSuggestions": [],
            "pages": [],
            "blocks": [],
            "sourceKind": "unknown",
            "error": str(exc),
        }
    sys.stdout.write(json.dumps(result, ensure_ascii=True))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
