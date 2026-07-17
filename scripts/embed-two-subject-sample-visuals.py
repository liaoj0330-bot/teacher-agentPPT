from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches


ROOT = Path(__file__).resolve().parents[1]
SOURCE_ROOT = ROOT / "artifacts" / "teacher-full-lesson-samples" / "browser-1784291858655"
VISUAL_ROOT = ROOT / "artifacts" / "teacher-full-lesson-samples" / "visual-assets-1784294700000"
OUTPUT_ROOT = ROOT / "artifacts" / "teacher-full-lesson-samples" / "final-visual"


def add_cover_crop(slide, image_path: Path, box: tuple[float, float, float, float]):
    x, y, width, height = box
    with Image.open(image_path) as image:
        image_width, image_height = image.size
    image_ratio = image_width / image_height
    box_ratio = width / height
    picture = slide.shapes.add_picture(
        str(image_path), Inches(x), Inches(y), width=Inches(width), height=Inches(height)
    )
    if image_ratio > box_ratio:
        visible = box_ratio / image_ratio
        picture.crop_left = picture.crop_right = (1 - visible) / 2
    else:
        visible = image_ratio / box_ratio
        picture.crop_top = picture.crop_bottom = (1 - visible) / 2
    picture.name = f"Sample visual - {image_path.stem}"
    return picture


def embed(source: Path, subject: str, placements: list[tuple[int, str, tuple[float, float, float, float]]]):
    presentation = Presentation(source)
    notes_before = [slide.notes_slide.notes_text_frame.text for slide in presentation.slides]
    for slide_number, asset_name, box in placements:
        add_cover_crop(presentation.slides[slide_number - 1], VISUAL_ROOT / subject / asset_name, box)

    target_dir = OUTPUT_ROOT / subject
    target_dir.mkdir(parents=True, exist_ok=True)
    target = target_dir / source.name
    presentation.save(target)

    reopened = Presentation(target)
    notes_after = [slide.notes_slide.notes_text_frame.text for slide in reopened.slides]
    image_count = sum(
        1
        for slide in reopened.slides
        for shape in slide.shapes
        if shape.name.startswith("Sample visual - ")
    )
    if len(reopened.slides) != 16 or image_count != len(placements) or notes_before != notes_after:
        raise RuntimeError(
            f"verification failed for {target.name}: slides={len(reopened.slides)}, images={image_count}, notes_preserved={notes_before == notes_after}"
        )
    return target, image_count


decks = list(SOURCE_ROOT.rglob("*.pptx"))
physics = next(path for path in decks if "物理" in path.name)
chinese = next(path for path in decks if "语文" in path.name)

outputs = [
    embed(
        physics,
        "physics",
        [
            (1, "physics-cover-lab.png", (7.30, 0.64, 5.31, 6.14)),
            (3, "physics-apparatus-closeup.png", (0.90, 1.67, 5.04, 4.61)),
            (10, "physics-student-inquiry.png", (8.67, 1.67, 3.74, 4.61)),
            (15, "physics-transfer-braking.png", (0.90, 1.67, 5.04, 4.61)),
        ],
    ),
    embed(
        chinese,
        "chinese",
        [
            (1, "chinese-cover-station.png", (7.30, 0.64, 5.31, 6.14)),
            (8, "chinese-father-platform.png", (7.30, 0.64, 5.31, 6.14)),
            (11, "chinese-detail-oranges.png", (8.67, 1.67, 3.74, 4.61)),
            (16, "chinese-life-transfer.png", (8.80, 1.67, 3.61, 4.61)),
        ],
    ),
]

for target, image_count in outputs:
    print(f"{target}\t16 slides\t{image_count} visuals")
