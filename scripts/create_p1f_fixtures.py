#!/usr/bin/env python
# -*- coding: utf-8 -*-

from __future__ import annotations

import sys
from pathlib import Path

from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
from reportlab.lib.pagesizes import A4
from reportlab.pdfbase.pdfmetrics import stringWidth
from reportlab.pdfgen import canvas


def create_docx(path: Path) -> None:
    doc = Document()
    doc.add_heading("universal-ppt-agent", level=1)
    doc.add_paragraph("evidence-regression source-document editable-pptx")
    doc.add_paragraph("This DOCX fixture checks real text extraction and paragraph parsing.")
    doc.add_paragraph("The parser should keep headings and visible text, not only the file name.")
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Column A"
    table.cell(0, 1).text = "Column B"
    table.cell(1, 0).text = "source-document"
    table.cell(1, 1).text = "editable-pptx"
    doc.save(str(path))


def create_pptx(path: Path) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    shapes = slide.shapes
    title = shapes.title or shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8.5), Inches(0.8))
    if shapes.title:
        title.text = "universal-ppt-agent"
    else:
        tf = title.text_frame
        tf.text = "universal-ppt-agent"
    textbox = shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(7.6), Inches(2.6))
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = "evidence-regression"
    p.font.size = Pt(20)
    p = tf.add_paragraph()
    p.text = "source-document editable-pptx"
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = "The parser should extract slide text and keep it traceable."
    p.font.size = Pt(16)
    prs.save(str(path))


def wrap_text(text: str, max_width: int, font_size: int = 11) -> list[str]:
    lines: list[str] = []
    current = ""
    for token in text.split():
        candidate = token if not current else f"{current} {token}"
        if stringWidth(candidate, "Helvetica", font_size) <= max_width:
            current = candidate
        else:
            if current:
                lines.append(current)
            current = token
    if current:
        lines.append(current)
    return lines


def create_pdf(path: Path) -> None:
    c = canvas.Canvas(str(path), pagesize=A4)
    width, height = A4
    c.setFont("Helvetica-Bold", 18)
    c.drawString(40, height - 60, "universal-ppt-agent")
    c.setFont("Helvetica", 12)
    text = "evidence-regression source-document editable-pptx"
    for index, line in enumerate(wrap_text(text, 500, 12)):
        c.drawString(40, height - 100 - index * 18, line)
    body = [
        "This PDF fixture checks text extraction without OCR.",
        "The parser should capture visible text and create stable page blocks.",
        "It is intentionally generic and reusable."
    ]
    y = height - 160
    for paragraph in body:
        for line in wrap_text(paragraph, 500, 11):
            c.drawString(40, y, line)
            y -= 16
        y -= 8
    c.showPage()
    c.save()


def main() -> int:
    if len(sys.argv) != 4:
      print("usage: create_p1f_fixtures.py docx pptx pdf", file=sys.stderr)
      return 2
    create_docx(Path(sys.argv[1]))
    create_pptx(Path(sys.argv[2]))
    create_pdf(Path(sys.argv[3]))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
